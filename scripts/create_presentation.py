"""
Generate PowerPoint presentation for Odin for Azure Local
Enhanced glassmorphism design matching the web application
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

# Color scheme matching the web app - enhanced for glassmorphism
COLORS = {
    'bg_dark': RGBColor(0x08, 0x08, 0x10),  # Deep dark with blue tint
    'card_bg': RGBColor(0x16, 0x16, 0x20),  # Glass card background
    'card_bg_light': RGBColor(0x20, 0x20, 0x2C),  # Lighter glass
    'text_primary': RGBColor(0xFF, 0xFF, 0xFF),
    'text_secondary': RGBColor(0xA8, 0xA8, 0xB3),
    'accent_blue': RGBColor(0x00, 0x78, 0xD4),
    'accent_blue_light': RGBColor(0x38, 0xB6, 0xFF),
    'accent_purple': RGBColor(0x8B, 0x5C, 0xF6),
    'accent_purple_light': RGBColor(0xA7, 0x8B, 0xFA),
    'success': RGBColor(0x10, 0xB9, 0x81),
    'glass_border': RGBColor(0x3C, 0x3C, 0x4A),
    'glass_border_light': RGBColor(0x50, 0x50, 0x60),
}

def set_slide_background(slide, color):
    """Set slide background to solid color"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_slide_decorations(slide):
    """Add decorative gradient orbs for glassmorphism effect"""
    # Top-right blue orb
    orb1 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9.5), Inches(-2), Inches(6), Inches(6))
    orb1.fill.solid()
    orb1.fill.fore_color.rgb = RGBColor(0x00, 0x28, 0x45)
    orb1.line.fill.background()
    
    # Bottom-left purple orb
    orb2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-2.5), Inches(4), Inches(5), Inches(5))
    orb2.fill.solid()
    orb2.fill.fore_color.rgb = RGBColor(0x25, 0x18, 0x40)
    orb2.line.fill.background()

def add_title_text(slide, text, left, top, width, height, font_size=44, color=None, bold=True, alignment=PP_ALIGN.CENTER):
    """Add a title text box"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color or COLORS['text_primary']
    p.alignment = alignment
    return txBox

def add_body_text(slide, text, left, top, width, height, font_size=18, color=None, alignment=PP_ALIGN.CENTER):
    """Add body text"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color or COLORS['text_secondary']
    p.alignment = alignment
    return txBox

def add_glass_card(slide, left, top, width, height, title, description, icon_text="", accent_color=None):
    """Add a feature card with enhanced glass-morphism style"""
    accent = accent_color or COLORS['accent_blue']
    
    # Outer glow effect
    glow = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, 
                                   left - Inches(0.02), top - Inches(0.02), 
                                   width + Inches(0.04), height + Inches(0.04))
    glow.fill.solid()
    glow.fill.fore_color.rgb = RGBColor(0x12, 0x12, 0x1A)
    glow.line.fill.background()
    
    # Card background with glass effect
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLORS['card_bg']
    shape.line.color.rgb = COLORS['glass_border']
    shape.line.width = Pt(1.5)
    
    # Top accent line
    accent_line = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, 
                                          left + Inches(0.15), top + Inches(0.12), 
                                          Inches(0.5), Inches(0.05))
    accent_line.fill.solid()
    accent_line.fill.fore_color.rgb = accent
    accent_line.line.fill.background()
    
    # Icon with glass background
    if icon_text:
        icon_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, 
                                          left + Inches(0.15), top + Inches(0.28), 
                                          Inches(0.52), Inches(0.52))
        icon_bg.fill.solid()
        icon_bg.fill.fore_color.rgb = COLORS['card_bg_light']
        icon_bg.line.color.rgb = COLORS['glass_border']
        icon_bg.line.width = Pt(1)
        
        icon_box = slide.shapes.add_textbox(left + Inches(0.15), top + Inches(0.32), Inches(0.52), Inches(0.48))
        tf = icon_box.text_frame
        p = tf.paragraphs[0]
        p.text = icon_text
        p.font.size = Pt(20)
        p.alignment = PP_ALIGN.CENTER
    
    # Title
    title_top = top + Inches(0.22) if not icon_text else top + Inches(0.9)
    title_box = slide.shapes.add_textbox(left + Inches(0.15), title_top, width - Inches(0.3), Inches(0.35))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = COLORS['text_primary']
    
    # Description
    desc_box = slide.shapes.add_textbox(left + Inches(0.15), title_top + Inches(0.3), width - Inches(0.3), height - Inches(1.1))
    tf = desc_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = description
    p.font.size = Pt(10)
    p.font.color.rgb = COLORS['text_secondary']

def add_step_badge(slide, left, top, number, label):
    """Add a step indicator badge with glass effect"""
    badge_width = Inches(1.5)
    badge_height = Inches(0.5)
    
    # Badge glass background
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, badge_width, badge_height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLORS['card_bg']
    shape.line.color.rgb = COLORS['glass_border']
    shape.line.width = Pt(1.5)
    
    # Number circle with accent
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, left + Inches(0.08), top + Inches(0.08), Inches(0.34), Inches(0.34))
    circle.fill.solid()
    circle.fill.fore_color.rgb = COLORS['accent_blue']
    circle.line.fill.background()
    
    # Number text
    num_box = slide.shapes.add_textbox(left + Inches(0.08), top + Inches(0.1), Inches(0.34), Inches(0.34))
    tf = num_box.text_frame
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    p = tf.paragraphs[0]
    p.text = str(number)
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = COLORS['text_primary']
    
    # Label
    label_box = slide.shapes.add_textbox(left + Inches(0.46), top + Inches(0.12), Inches(1), Inches(0.3))
    tf = label_box.text_frame
    p = tf.paragraphs[0]
    p.text = label
    p.font.size = Pt(11)
    p.font.color.rgb = COLORS['text_primary']

def add_arrow(slide, left, top):
    """Add arrow between steps"""
    arrow_box = slide.shapes.add_textbox(left, top, Inches(0.3), Inches(0.5))
    tf = arrow_box.text_frame
    p = tf.paragraphs[0]
    p.text = "→"
    p.font.size = Pt(18)
    p.font.color.rgb = COLORS['accent_blue_light']
    p.alignment = PP_ALIGN.CENTER

def add_highlight_box(slide, left, top, width, height, title, text):
    """Add a highlighted info box with glow effect"""
    # Outer glow
    glow = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, 
                                   left - Inches(0.04), top - Inches(0.04), 
                                   width + Inches(0.08), height + Inches(0.08))
    glow.fill.solid()
    glow.fill.fore_color.rgb = RGBColor(0x00, 0x30, 0x55)
    glow.line.fill.background()
    
    # Main box
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x0C, 0x1C, 0x30)
    shape.line.color.rgb = COLORS['accent_blue']
    shape.line.width = Pt(1.5)
    
    # Title
    title_box = slide.shapes.add_textbox(left + Inches(0.25), top + Inches(0.12), width - Inches(0.5), Inches(0.35))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = COLORS['accent_blue_light']
    
    # Text
    text_box = slide.shapes.add_textbox(left + Inches(0.25), top + Inches(0.45), width - Inches(0.5), height - Inches(0.55))
    tf = text_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(12)
    p.font.color.rgb = COLORS['text_secondary']

def add_stat_card(slide, left, top, number, label):
    """Add a statistic display with glass panel"""
    # Outer glow
    glow = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, 
                                   left - Inches(0.02), top - Inches(0.02), 
                                   Inches(2.24), Inches(1.44))
    glow.fill.solid()
    glow.fill.fore_color.rgb = RGBColor(0x10, 0x10, 0x18)
    glow.line.fill.background()
    
    # Glass panel
    panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(2.2), Inches(1.4))
    panel.fill.solid()
    panel.fill.fore_color.rgb = COLORS['card_bg']
    panel.line.color.rgb = COLORS['glass_border']
    panel.line.width = Pt(1.5)
    
    # Number
    num_box = slide.shapes.add_textbox(left, top + Inches(0.15), Inches(2.2), Inches(0.7))
    tf = num_box.text_frame
    p = tf.paragraphs[0]
    p.text = number
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = COLORS['accent_blue_light']
    p.alignment = PP_ALIGN.CENTER
    
    # Label
    label_box = slide.shapes.add_textbox(left, top + Inches(0.9), Inches(2.2), Inches(0.4))
    tf = label_box.text_frame
    p = tf.paragraphs[0]
    p.text = label
    p.font.size = Pt(11)
    p.font.color.rgb = COLORS['text_secondary']
    p.alignment = PP_ALIGN.CENTER

def add_bullet_list(slide, left, top, width, items, font_size=14):
    """Add a bullet list with checkmarks and glass backgrounds"""
    current_top = top
    for item in items:
        # Glass pill background for each item
        pill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, 
                                       left - Inches(0.1), current_top - Inches(0.02), 
                                       width + Inches(0.2), Inches(0.42))
        pill.fill.solid()
        pill.fill.fore_color.rgb = COLORS['card_bg']
        pill.line.color.rgb = COLORS['glass_border']
        pill.line.width = Pt(1)
        
        # Checkmark circle
        check_bg = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, current_top + Inches(0.04), Inches(0.3), Inches(0.3))
        check_bg.fill.solid()
        check_bg.fill.fore_color.rgb = RGBColor(0x0D, 0x35, 0x28)
        check_bg.line.color.rgb = COLORS['success']
        check_bg.line.width = Pt(1)
        
        check_box = slide.shapes.add_textbox(left, current_top + Inches(0.02), Inches(0.3), Inches(0.34))
        tf = check_box.text_frame
        p = tf.paragraphs[0]
        p.text = "✓"
        p.font.size = Pt(font_size)
        p.font.color.rgb = COLORS['success']
        p.alignment = PP_ALIGN.CENTER
        
        # Text
        text_box = slide.shapes.add_textbox(left + Inches(0.4), current_top + Inches(0.06), width - Inches(0.4), Inches(0.38))
        tf = text_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = COLORS['text_primary']
        
        current_top += Inches(0.48)

def add_value_card(slide, left, top, width, height, icon, title, description):
    """Add a value proposition card with enhanced glass styling"""
    # Outer glow
    glow = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, 
                                   left - Inches(0.03), top - Inches(0.03), 
                                   width + Inches(0.06), height + Inches(0.06))
    glow.fill.solid()
    glow.fill.fore_color.rgb = RGBColor(0x12, 0x12, 0x1C)
    glow.line.fill.background()
    
    # Card background
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLORS['card_bg']
    shape.line.color.rgb = COLORS['glass_border_light']
    shape.line.width = Pt(2)
    
    # Top accent line
    accent = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, 
                                     left + Inches(0.15), top + Inches(0.15), 
                                     width - Inches(0.3), Inches(0.05))
    accent.fill.solid()
    accent.fill.fore_color.rgb = COLORS['accent_blue']
    accent.line.fill.background()
    
    # Icon background circle with border
    icon_size = Inches(0.75)
    icon_left = left + (width - icon_size) / 2
    icon_bg = slide.shapes.add_shape(MSO_SHAPE.OVAL, icon_left, top + Inches(0.35), icon_size, icon_size)
    icon_bg.fill.solid()
    icon_bg.fill.fore_color.rgb = COLORS['card_bg_light']
    icon_bg.line.color.rgb = COLORS['accent_blue']
    icon_bg.line.width = Pt(1.5)
    
    # Icon
    icon_box = slide.shapes.add_textbox(icon_left, top + Inches(0.45), icon_size, icon_size)
    tf = icon_box.text_frame
    p = tf.paragraphs[0]
    p.text = icon
    p.font.size = Pt(28)
    p.alignment = PP_ALIGN.CENTER
    
    # Title
    title_box = slide.shapes.add_textbox(left, top + Inches(1.2), width, Inches(0.4))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = COLORS['accent_blue_light']
    p.alignment = PP_ALIGN.CENTER
    
    # Description
    desc_box = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(1.6), width - Inches(0.4), height - Inches(1.8))
    tf = desc_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = description
    p.font.size = Pt(11)
    p.font.color.rgb = COLORS['text_secondary']
    p.alignment = PP_ALIGN.CENTER

def create_presentation():
    """Create the complete PowerPoint presentation"""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    blank_layout = prs.slide_layouts[6]
    
    # Get logo path
    script_dir = os.path.dirname(os.path.abspath(__file__))
    logo_path = os.path.join(script_dir, "images", "odin-logo.png")
    
    # =========== SLIDE 1: Title ===========
    slide1 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide1, COLORS['bg_dark'])
    add_slide_decorations(slide1)
    
    # Add logo image
    if os.path.exists(logo_path):
        logo = slide1.shapes.add_picture(logo_path, Inches(5.15), Inches(0.5), width=Inches(3))
    
    # Main title
    add_title_text(slide1, "Odin for Azure Local", Inches(0), Inches(3.0), Inches(13.333), Inches(0.8), 
                   font_size=52, color=COLORS['text_primary'])
    
    # Subtitle with accent color
    subtitle = slide1.shapes.add_textbox(Inches(0), Inches(3.8), Inches(13.333), Inches(0.5))
    tf = subtitle.text_frame
    p = tf.paragraphs[0]
    p.text = "Optimal Deployment and Infrastructure Navigator"
    p.font.size = Pt(22)
    p.font.color.rgb = COLORS['accent_blue_light']
    p.alignment = PP_ALIGN.CENTER
    
    # Description
    add_body_text(slide1, "A comprehensive web-based wizard to design, configure, and deploy\nAzure Local (formerly Azure Stack HCI) network architecture with confidence.",
                  Inches(1.5), Inches(4.4), Inches(10.333), Inches(0.8), font_size=14, color=COLORS['text_secondary'])
    
    # Stats with glass panels
    add_stat_card(slide1, Inches(2.6), Inches(5.4), "20+", "Configuration Steps")
    add_stat_card(slide1, Inches(5.5), Inches(5.4), "4", "Deployment Scenarios")
    add_stat_card(slide1, Inches(8.4), Inches(5.4), "100%", "Client-Side")
    
    # =========== SLIDE 2: What is Odin? ===========
    slide2 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide2, COLORS['bg_dark'])
    add_slide_decorations(slide2)
    
    add_title_text(slide2, "What is Odin?", Inches(0), Inches(0.4), Inches(13.333), Inches(0.7), font_size=40)
    add_body_text(slide2, "Named after the Norse god of wisdom and strategy", 
                  Inches(0), Inches(1.0), Inches(13.333), Inches(0.4), font_size=16, color=COLORS['accent_purple_light'])
    
    add_body_text(slide2, "Odin is a guided decision-tree interface that helps IT architects and engineers\ndesign Azure Local cluster configurations using validated architecture patterns.",
                  Inches(1.5), Inches(1.5), Inches(10.333), Inches(0.7), font_size=14, color=COLORS['text_secondary'])
    
    # Highlight box
    add_highlight_box(slide2, Inches(3.2), Inches(2.3), Inches(6.9), Inches(0.85),
                      "🎯 Core Purpose",
                      "Accelerate skills ramp-up for Azure Local while helping validate cluster design configurations.")
    
    # Feature list
    features = [
        "Step-by-Step Wizard — Guided flow through all deployment decisions",
        "Real-Time Validation — Instant feedback on all configuration inputs",
        "ARM Template Generation — Export ready-to-deploy Azure parameters",
        "Import Existing Deployments — Generate docs from deployed ARM templates",
        "Visual Architecture Diagrams — See your network topology in real-time"
    ]
    add_bullet_list(slide2, Inches(2.3), Inches(3.4), Inches(8.7), features, font_size=12)
    
    # =========== SLIDE 3: Deployment Scenarios ===========
    slide3 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide3, COLORS['bg_dark'])
    add_slide_decorations(slide3)
    
    add_title_text(slide3, "Supported Deployment Scenarios", Inches(0), Inches(0.4), Inches(13.333), Inches(0.7), font_size=40)
    add_body_text(slide3, "Choose the right architecture for your needs", 
                  Inches(0), Inches(1.0), Inches(13.333), Inches(0.4), font_size=16, color=COLORS['text_secondary'])
    
    # Scenario cards
    card_width = Inches(2.9)
    card_height = Inches(2.5)
    start_left = Inches(0.75)
    card_top = Inches(1.65)
    gap = Inches(0.25)
    
    scenarios = [
        ("🖥️", "Hyperconverged", "Single or rack-aware cluster with compute, storage & network combined", COLORS['accent_blue']),
        ("🏢", "Multi-Rack", "Scalable distributed architecture across multiple racks", COLORS['accent_purple']),
        ("🛡️", "Disconnected", "Air-gapped operation with local management", COLORS['accent_blue']),
        ("📊", "M365 Local", "Microsoft 365 workloads with minimum 9 nodes", COLORS['accent_purple'])
    ]
    
    for i, (icon, title, desc, accent) in enumerate(scenarios):
        left = start_left + (card_width + gap) * i
        add_glass_card(slide3, left, card_top, card_width, card_height, title, desc, icon, accent)
    
    # Additional info
    add_highlight_box(slide3, Inches(2.2), Inches(4.5), Inches(8.9), Inches(0.85),
                      "📘 Documentation Links",
                      "Each scenario includes links to official Microsoft documentation for detailed requirements.")
    
    # =========== SLIDE 4: Core Features ===========
    slide4 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide4, COLORS['bg_dark'])
    add_slide_decorations(slide4)
    
    add_title_text(slide4, "Core Features", Inches(0), Inches(0.3), Inches(13.333), Inches(0.6), font_size=40)
    add_body_text(slide4, "Everything you need to design Azure Local deployments", 
                  Inches(0), Inches(0.85), Inches(13.333), Inches(0.35), font_size=14, color=COLORS['text_secondary'])
    
    # Feature cards in 2x3 grid
    features = [
        ("📄", "ARM Template Generation", "Generate Azure Resource Manager parameters JSON ready for deployment.", COLORS['accent_blue']),
        ("📥", "Import Deployed Templates", "Import ARM templates from existing deployments to generate documentation.", COLORS['accent_purple']),
        ("📊", "Visual Diagrams", "Real-time network topology and architecture visualizations.", COLORS['accent_blue']),
        ("📋", "Documentation Export", "Generate comprehensive reports (HTML/Word) from any configuration.", COLORS['accent_purple']),
        ("✓", "Real-Time Validation", "Instant feedback on all inputs with helpful error messages.", COLORS['accent_blue']),
        ("💾", "Auto-Save & Resume", "Progress automatically saved. Resume anytime from where you left off.", COLORS['accent_purple'])
    ]
    
    card_width = Inches(4)
    card_height = Inches(1.7)
    start_left = Inches(0.5)
    start_top = Inches(1.35)
    h_gap = Inches(0.25)
    v_gap = Inches(0.2)
    
    for i, (icon, title, desc, accent) in enumerate(features):
        row = i // 3
        col = i % 3
        left = start_left + (card_width + h_gap) * col
        top = start_top + (card_height + v_gap) * row
        add_glass_card(slide4, left, top, card_width, card_height, title, desc, icon, accent)
    
    # =========== SLIDE 5: Wizard Flow ===========
    slide5 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide5, COLORS['bg_dark'])
    add_slide_decorations(slide5)
    
    add_title_text(slide5, "Guided Configuration Flow", Inches(0), Inches(0.4), Inches(13.333), Inches(0.7), font_size=40)
    add_body_text(slide5, "Logical step-by-step progression through all decisions", 
                  Inches(0), Inches(1.0), Inches(13.333), Inches(0.4), font_size=16, color=COLORS['text_secondary'])
    
    # First row of steps
    steps_row1 = ["Type", "Cloud", "Region", "Scale"]
    start_left = Inches(1.3)
    step_top = Inches(1.75)
    
    for i, label in enumerate(steps_row1):
        left = start_left + Inches(2.8) * i
        add_step_badge(slide5, left, step_top, i + 1, label)
        if i < len(steps_row1) - 1:
            add_arrow(slide5, left + Inches(1.6), step_top)
    
    # Second row of steps
    steps_row2 = ["Nodes", "Infrastructure", "Identity", "Security"]
    step_top2 = Inches(2.5)
    
    for i, label in enumerate(steps_row2):
        left = start_left + Inches(2.8) * i
        add_step_badge(slide5, left, step_top2, i + 5, label)
        if i < len(steps_row2) - 1:
            add_arrow(slide5, left + Inches(1.6), step_top2)
    
    # Built-in tools highlight
    add_highlight_box(slide5, Inches(2.2), Inches(3.5), Inches(8.9), Inches(0.85),
                      "🔧 Built-in Tools",
                      "CIDR Calculator  •  Cost Estimator  •  Configuration Preview  •  Report Generator")
    
    # Tool cards
    tools = [
        ("📐", "CIDR Calculator", "Calculate subnet details, usable IPs, and network ranges", COLORS['accent_blue']),
        ("💰", "Cost Estimator", "Rough monthly cost estimates based on your configuration", COLORS['accent_purple']),
        ("📋", "Report Generator", "Export comprehensive HTML/Word reports with diagrams", COLORS['accent_blue'])
    ]
    
    for i, (icon, title, desc, accent) in enumerate(tools):
        left = Inches(0.8) + Inches(4.1) * i
        add_glass_card(slide5, left, Inches(4.7), Inches(3.9), Inches(1.6), title, desc, icon, accent)
    
    # =========== SLIDE 6: Value Proposition ===========
    slide6 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide6, COLORS['bg_dark'])
    add_slide_decorations(slide6)
    
    add_title_text(slide6, "Value Proposition", Inches(0), Inches(0.4), Inches(13.333), Inches(0.7), font_size=40)
    add_body_text(slide6, "Why use Odin for your Azure Local deployments?", 
                  Inches(0), Inches(1.0), Inches(13.333), Inches(0.4), font_size=16, color=COLORS['text_secondary'])
    
    # Value cards
    values = [
        ("⚡", "Accelerate Learning", "Quickly understand Azure Local concepts, network requirements, and deployment options."),
        ("📚", "Document Existing", "Import ARM templates from deployed clusters to generate architecture documentation."),
        ("🚀", "Deploy Faster", "Generate ARM parameters directly from your design. One-click deployment to Azure Portal.")
    ]
    
    card_width = Inches(3.9)
    card_height = Inches(2.6)
    start_left = Inches(0.8)
    
    for i, (icon, title, desc) in enumerate(values):
        left = start_left + (card_width + Inches(0.3)) * i
        add_value_card(slide6, left, Inches(1.55), card_width, card_height, icon, title, desc)
    
    # Bottom highlight
    add_highlight_box(slide6, Inches(2.7), Inches(4.5), Inches(7.9), Inches(0.9),
                      "💡 No Installation Required",
                      "100% client-side web application. Works offline. No data leaves your browser.")
    
    # =========== SLIDE 7: Technical Highlights ===========
    slide7 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide7, COLORS['bg_dark'])
    add_slide_decorations(slide7)
    
    add_title_text(slide7, "Technical Highlights", Inches(0), Inches(0.4), Inches(13.333), Inches(0.7), font_size=40)
    add_body_text(slide7, "Built with simplicity and security in mind", 
                  Inches(0), Inches(1.0), Inches(13.333), Inches(0.4), font_size=16, color=COLORS['text_secondary'])
    
    # Technical feature cards
    tech_features = [
        ("🛡️", "Client-Side Only", "No server, no data collection. Everything runs in your browser with localStorage.", COLORS['accent_blue']),
        ("📦", "No Dependencies", "Pure HTML, CSS, and vanilla JavaScript. No frameworks, no build step required.", COLORS['accent_purple']),
        ("🌐", "Works Offline", "Once loaded, works without internet. Perfect for air-gapped planning scenarios.", COLORS['accent_blue']),
        ("☁️", "ARM Integration", "Import existing ARM templates. Export ready-to-deploy parameters for Azure.", COLORS['accent_purple'])
    ]
    
    card_width = Inches(6.2)
    card_height = Inches(1.5)
    start_left = Inches(0.4)
    start_top = Inches(1.5)
    h_gap = Inches(0.25)
    v_gap = Inches(0.2)
    
    for i, (icon, title, desc, accent) in enumerate(tech_features):
        row = i // 2
        col = i % 2
        left = start_left + (card_width + h_gap) * col
        top = start_top + (card_height + v_gap) * row
        add_glass_card(slide7, left, top, card_width, card_height, title, desc, icon, accent)
    
    # Technology stack with glass pills
    stack_top = Inches(4.75)
    add_body_text(slide7, "Technology Stack", Inches(0), stack_top, Inches(13.333), Inches(0.4), 
                  font_size=14, color=COLORS['accent_blue_light'])
    
    stack_items = ["HTML5", "CSS3", "Vanilla JavaScript", "localStorage API"]
    item_width = Inches(2.6)
    start_x = Inches(1.8)
    
    for i, item in enumerate(stack_items):
        left = start_x + item_width * i
        # Glass pill
        pill = slide7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, stack_top + Inches(0.45), Inches(2.4), Inches(0.5))
        pill.fill.solid()
        pill.fill.fore_color.rgb = COLORS['card_bg']
        pill.line.color.rgb = COLORS['accent_blue']
        pill.line.width = Pt(1.5)
        
        text_box = slide7.shapes.add_textbox(left, stack_top + Inches(0.53), Inches(2.4), Inches(0.4))
        tf = text_box.text_frame
        p = tf.paragraphs[0]
        p.text = item
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = COLORS['text_primary']
        p.alignment = PP_ALIGN.CENTER
    
    # =========== SLIDE 8: Get Started ===========
    slide8 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide8, COLORS['bg_dark'])
    add_slide_decorations(slide8)
    
    add_title_text(slide8, "Get Started Today", Inches(0), Inches(0.4), Inches(13.333), Inches(0.7), font_size=40)
    add_body_text(slide8, "Begin designing your Azure Local deployment in minutes", 
                  Inches(0), Inches(1.0), Inches(13.333), Inches(0.4), font_size=16, color=COLORS['text_secondary'])
    
    # Quick start box
    add_highlight_box(slide8, Inches(2.2), Inches(1.55), Inches(8.9), Inches(0.9),
                      "📖 Quick Start",
                      "Open index.html in a browser, follow the wizard steps, export your ARM parameters, and deploy to Azure!")
    
    # Steps with glass panels
    steps = [
        ("1", "Open the Wizard", "Launch index.html or use serve.ps1 for local development"),
        ("2", "Answer Questions", "Follow the guided flow to configure your deployment"),
        ("3", "Review & Validate", "Check the real-time summary and diagrams"),
        ("4", "Export & Deploy", "Generate ARM parameters and deploy to Azure Portal")
    ]
    
    current_top = Inches(2.75)
    for num, title, desc in steps:
        # Glass panel for each step
        panel = slide8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2.2), current_top - Inches(0.05), Inches(8.9), Inches(0.52))
        panel.fill.solid()
        panel.fill.fore_color.rgb = COLORS['card_bg']
        panel.line.color.rgb = COLORS['glass_border']
        panel.line.width = Pt(1.5)
        
        # Number circle
        circle = slide8.shapes.add_shape(MSO_SHAPE.OVAL, Inches(2.35), current_top, Inches(0.42), Inches(0.42))
        circle.fill.solid()
        circle.fill.fore_color.rgb = COLORS['success']
        circle.line.fill.background()
        
        num_box = slide8.shapes.add_textbox(Inches(2.35), current_top + Inches(0.06), Inches(0.42), Inches(0.36))
        tf = num_box.text_frame
        p = tf.paragraphs[0]
        p.text = num
        p.font.size = Pt(15)
        p.font.bold = True
        p.font.color.rgb = COLORS['text_primary']
        p.alignment = PP_ALIGN.CENTER
        
        # Title and description
        text_box = slide8.shapes.add_textbox(Inches(2.9), current_top + Inches(0.1), Inches(8), Inches(0.4))
        tf = text_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"{title} — {desc}"
        p.font.size = Pt(13)
        p.font.color.rgb = COLORS['text_primary']
        
        current_top += Inches(0.58)
    
    # CTA Button with glow
    cta_glow = slide8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.2), Inches(5.2), Inches(4.9), Inches(0.85))
    cta_glow.fill.solid()
    cta_glow.fill.fore_color.rgb = RGBColor(0x00, 0x3D, 0x6A)
    cta_glow.line.fill.background()
    
    cta_shape = slide8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.3), Inches(5.3), Inches(4.7), Inches(0.65))
    cta_shape.fill.solid()
    cta_shape.fill.fore_color.rgb = COLORS['accent_blue']
    cta_shape.line.fill.background()
    
    cta_text = slide8.shapes.add_textbox(Inches(4.3), Inches(5.4), Inches(4.7), Inches(0.55))
    tf = cta_text.text_frame
    p = tf.paragraphs[0]
    p.text = "▶  Launch Odin Wizard"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = COLORS['text_primary']
    p.alignment = PP_ALIGN.CENTER
    
    # Footer
    footer = slide8.shapes.add_textbox(Inches(0), Inches(6.7), Inches(13.333), Inches(0.4))
    tf = footer.text_frame
    p = tf.paragraphs[0]
    p.text = "Open index.html in your browser to start  •  Odin for Azure Local"
    p.font.size = Pt(11)
    p.font.color.rgb = COLORS['text_secondary']
    p.alignment = PP_ALIGN.CENTER
    
    # Save presentation
    output_path = os.path.join(script_dir, "Odin_Presentation.pptx")
    prs.save(output_path)
    print(f"✅ Presentation saved to: {output_path}")
    return output_path

if __name__ == "__main__":
    create_presentation()
